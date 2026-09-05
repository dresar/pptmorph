import os
import copy
import zipfile
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import xml.etree.ElementTree as ET

NS = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
}

TARGET_PPTX = r"C:\Users\NCN0C\Downloads\Desain tanpa judul.pptx"

def get_shape_transform(elem):
    off = elem.find('.//a:off', NS)
    ext = elem.find('.//a:ext', NS)
    x = int(off.get('x', 0)) if off is not None else 0
    y = int(off.get('y', 0)) if off is not None else 0
    cx = int(ext.get('cx', 0)) if ext is not None else 0
    cy = int(ext.get('cy', 0)) if ext is not None else 0
    return x, y, cx, cy

def set_shape_position(elem, new_x, new_y):
    off = elem.find('.//a:off', NS)
    if off is not None:
        off.set('x', str(int(new_x)))
        off.set('y', str(int(new_y)))

def set_shape_id_name(elem, shape_id, shape_name):
    cnvpr = elem.find('.//p:cNvPr', NS)
    if cnvpr is not None:
        cnvpr.set('id', str(shape_id))
        cnvpr.set('name', shape_name)

def build_master_presentation():
    # 1. Load the presentation
    prs = Presentation(TARGET_PPTX)
    w = prs.slide_width
    h = prs.slide_height

    # Extract all original in-canvas shapes from the 12 slides
    # (Extract only authentic in-canvas shapes, ignoring previously generated off-canvas clones)
    original_slides_data = [] # List of list of shape info dicts
    
    for s_idx, slide in enumerate(prs.slides):
        slide_shapes = []
        for s in list(slide.shapes):
            x, y, cx, cy = get_shape_transform(s._element)
            txt = s.text_frame.text.replace('\n', ' ').strip() if s.has_text_frame else ''
            
            # Check if authentic in-canvas shape
            in_canvas = (s.left >= -200000 and s.top >= -200000 and s.left + s.width <= w + 500000 and s.top + s.height <= h + 500000)
            
            # Persistent background covers entire slide (left=0, top=0, width=w, height=h)
            is_bg = (s.left == 0 and s.top == 0 and s.width >= w and s.height >= h)
            
            if in_canvas or is_bg:
                # Assign a canonical name
                canonical_name = f"!!Slide{s_idx+1}_Shp{s.shape_id}"
                s.name = canonical_name
                
                elem_copy = copy.deepcopy(s._element)
                slide_shapes.append({
                    'elem': elem_copy,
                    'orig_x': x,
                    'orig_y': y,
                    'cx': cx,
                    'cy': cy,
                    'name': canonical_name,
                    'is_bg': is_bg,
                    'text': txt
                })
        original_slides_data.append(slide_shapes)

    # 2. Re-create a clean presentation with Slide 0 (White intro) + 12 slides
    clean_prs = Presentation(TARGET_PPTX)
    
    # Add blank slide for Slide 0 (Intro White Screen)
    blank_layout = clean_prs.slide_layouts[6]
    intro_slide = clean_prs.slides.add_slide(blank_layout)
    
    # Move intro_slide to index 0
    sldIdLst = clean_prs._element.sldIdLst
    last_id = sldIdLst[-1]
    sldIdLst.remove(last_id)
    sldIdLst.insert(0, last_id)
    
    # Add a white background shape to intro_slide so it is 100% crisp pure white
    white_bg = intro_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, w, h
    )
    white_bg.fill.solid()
    white_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    white_bg.line.fill.background()
    white_bg.name = "!!Intro_White_Canvas"

    # Save to temp structure
    temp_clean = r"C:\Users\NCN0C\Pictures\PPTMORPH\temp_clean.pptx"
    clean_prs.save(temp_clean)

    # 3. Re-open to apply complete dual-stage Off-Canvas Staging
    prs = Presentation(temp_clean)
    num_total_slides = len(prs.slides) # 13 slides (0 is intro, 1..12 are main)

    def get_max_id(slide):
        max_id = 200
        for sh in slide.shapes:
            if sh.shape_id > max_id:
                max_id = sh.shape_id
        return max_id + 20

    # Total slide data list:
    # Index 0: Slide 0 (Intro - contains only white_bg)
    # Index 1..12: Slide 1..12 (Original contents)
    all_slide_data = [
        [{
            'elem': copy.deepcopy(white_bg._element),
            'orig_x': 0, 'orig_y': 0, 'cx': w, 'cy': h,
            'name': "!!Intro_White_Canvas", 'is_bg': True, 'text': ''
        }]
    ] + original_slides_data

    # STAGING PIPELINE:
    for i in range(num_total_slides):
        curr_slide = prs.slides[i]
        sp_tree = curr_slide._element.spTree
        next_id = get_max_id(curr_slide)

        # -------------------------------------------------------------
        # A. INBOUND STAGING: Upcoming shapes from Slide i+1 into Slide i
        # -------------------------------------------------------------
        if i + 1 < num_total_slides:
            next_items = all_slide_data[i + 1]
            for idx_item, item in enumerate(next_items):
                elem = copy.deepcopy(item['elem'])
                orig_x = item['orig_x']
                orig_y = item['orig_y']
                cx = item['cx']
                cy = item['cy']
                is_bg = item['is_bg']
                stagger = (idx_item % 5) * 400000

                # Off-Canvas Entry Trajectories:
                if is_bg:
                    # Background enters by scaling or from top
                    entry_x = 0
                    entry_y = int(-h - 1000000)
                elif orig_y < h * 0.35:
                    # Top headers and titles enter from Top or Right
                    if idx_item % 2 == 0:
                        entry_x = orig_x
                        entry_y = int(-cy - 2500000 - stagger)
                    else:
                        entry_x = int(w + 2500000 + stagger)
                        entry_y = orig_y
                elif orig_y > h * 0.55:
                    # Bottom cards, metrics, contact info enter from Bottom
                    entry_x = orig_x
                    entry_y = int(h + 2500000 + stagger)
                else:
                    # Middle content enters from Right
                    entry_x = int(w + 3000000 + stagger)
                    entry_y = orig_y

                set_shape_position(elem, entry_x, entry_y)
                set_shape_id_name(elem, next_id, item['name'])
                next_id += 1
                sp_tree.append(elem)

        # -------------------------------------------------------------
        # B. OUTBOUND STAGING: Exiting shapes from Slide i-1 into Slide i
        # -------------------------------------------------------------
        if i > 0:
            prev_items = all_slide_data[i - 1]
            for idx_item, item in enumerate(prev_items):
                elem = copy.deepcopy(item['elem'])
                orig_x = item['orig_x']
                orig_y = item['orig_y']
                cx = item['cx']
                cy = item['cy']
                is_bg = item['is_bg']
                stagger = (idx_item % 5) * 400000

                # Off-Canvas Exit Trajectories:
                if is_bg:
                    # Previous background exits to top/left
                    exit_x = int(-w - 1000000)
                    exit_y = 0
                elif orig_y < h * 0.4:
                    # Top titles/headers exit to Top
                    exit_x = orig_x
                    exit_y = int(-cy - 2500000 - stagger)
                elif orig_y > h * 0.6:
                    # Bottom elements exit to Bottom or Left
                    if idx_item % 2 == 0:
                        exit_x = int(-cx - 3000000 - stagger)
                        exit_y = orig_y
                    else:
                        exit_x = orig_x
                        exit_y = int(h + 2500000 + stagger)
                else:
                    # Middle content exits to Left
                    exit_x = int(-cx - 3500000 - stagger)
                    exit_y = orig_y

                set_shape_position(elem, exit_x, exit_y)
                set_shape_id_name(elem, next_id, item['name'])
                next_id += 1
                sp_tree.append(elem)

    temp_staged = r"C:\Users\NCN0C\Pictures\PPTMORPH\temp_staged.pptx"
    prs.save(temp_staged)

    # 4. Inject Morph Transitions into all slides (including Slide 1 from Intro Slide 0)
    morph_xml = (
        '<mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006">'
        '<mc:Choice xmlns:p159="http://schemas.microsoft.com/office/powerpoint/2015/09/main" Requires="p159">'
        '<p:transition spd="slow" xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" p14:dur="1250">'
        '<p159:morph option="byObject"/>'
        '</p:transition>'
        '</mc:Choice>'
        '<mc:Fallback><p:transition spd="med"><p:fade/></p:transition></mc:Fallback>'
        '</mc:AlternateContent>'
    )

    temp_final = r"C:\Users\NCN0C\Pictures\PPTMORPH\temp_final.pptx"
    with zipfile.ZipFile(temp_staged, 'r') as zin, zipfile.ZipFile(temp_final, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith('ppt/slides/slide') and item.filename.endswith('.xml'):
                xml_str = data.decode('utf-8')
                
                # Clean any existing transition tags
                xml_str = re.sub(r'<p:transition.*?</p:transition>', '', xml_str)
                xml_str = re.sub(r'<mc:AlternateContent.*?</mc:AlternateContent>', '', xml_str)
                
                # All slides get Morph Transition!
                xml_str = xml_str.replace('</p:sld>', f'{morph_xml}</p:sld>')
                data = xml_str.encode('utf-8')
            zout.writestr(item, data)

    # Directly overwrite TARGET_PPTX
    if os.path.exists(TARGET_PPTX):
        os.remove(TARGET_PPTX)
    os.rename(temp_final, TARGET_PPTX)

    # Cleanup temporary working files
    for tmp in [temp_clean, temp_staged]:
        if os.path.exists(tmp):
            os.remove(tmp)

    print(f"Master Morph Presentation successfully built into: {TARGET_PPTX}")

if __name__ == "__main__":
    build_master_presentation()
