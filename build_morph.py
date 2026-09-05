import os
import shutil
import zipfile
import re
from pptx import Presentation

INPUT_PPTX = r"C:\Users\NCN0C\Downloads\Desain tanpa judul.pptx"
BACKUP_PPTX = r"C:\Users\NCN0C\Downloads\Desain tanpa judul_backup.pptx"
OUTPUT_DOWNLOADS_MORPH = r"C:\Users\NCN0C\Downloads\Desain tanpa judul_morph.pptx"
OUTPUT_WORKSPACE_MORPH = r"C:\Users\NCN0C\Pictures\PPTMORPH\Desain tanpa judul_morph.pptx"

def setup_morph_presentation():
    # 1. Create backup
    if not os.path.exists(BACKUP_PPTX):
        shutil.copy2(INPUT_PPTX, BACKUP_PPTX)
        print(f"Backup created: {BACKUP_PPTX}")

    prs = Presentation(INPUT_PPTX)
    
    # Title mappings for each slide
    slide_title_words = {
        0: {"h1": "PORTOFOLIO", "h2": "Kevin Pratama"},
        1: {"h1": "TENTANG", "h2": "SAYA"},
        2: {"h1": "PROFIL", "h2": "DIRI"},
        3: {"h1": "KEAHLIAN", "h2": "SAYA"},
        4: {"h1": "RIWAYAT", "h2": "PENDIDIKAN"},
        5: {"h1": "SERTIFIKASI", "h2": "SAYA"},
        6: {"h1": "BIDANG", "h2": "KONTEN"},
        7: {"h1": "PORTOFOLIO", "h2": "PILIHAN"},
        8: {"h1": "PENCAPAIAN", "h2": "KREATIF"},
        9: {"h1": "PENGALAMAN", "h2": "PROFESIONAL"},
        10: {"h1": "LAYANAN", "h2": "KREATIF"},
        11: {"h1": "TERIMA", "h2": "KASIH"},
    }

    for idx, slide in enumerate(prs.slides):
        mapping = slide_title_words.get(idx, {})
        h1_target = mapping.get("h1")
        h2_target = mapping.get("h2")
        
        for s in slide.shapes:
            text = s.text_frame.text.strip() if s.has_text_frame else ""
            
            # Synchronize Header Line
            if s.top in range(1000000, 1050000) and s.height < 20000 and s.width > 10000000:
                s.name = "!!HeaderLine"
            # Synchronize Header Branding / Subtitle
            elif "Portofolio - Kevin Pratama" in text:
                s.name = "!!HeaderTitle"
            # Synchronize Slide Number
            elif text in [f"{i:02d}" for i in range(1, 13)] or text in [f"{i}" for i in range(1, 13)]:
                s.name = "!!SlideNumber"
            # Synchronize Main Big Heading 1
            elif h1_target and (text == h1_target or text.startswith(h1_target)):
                s.name = "!!MainHeading1"
            # Synchronize Main Big Heading 2
            elif h2_target and (text == h2_target or text.startswith(h2_target)):
                s.name = "!!MainHeading2"
            # Synchronize Background shape
            elif s.shape_id == 2 and s.left == 0 and s.top == 0:
                s.name = "!!Background"
            # Synchronize Left Accent Decor
            elif s.name in ["Group 7", "Group 4"]:
                s.name = f"!!Decor_{s.name.replace(' ', '_')}"

    # Save to temp file
    temp_saved = r"C:\Users\NCN0C\Pictures\PPTMORPH\temp_renamed.pptx"
    prs.save(temp_saved)
    print("Shape names synchronized with Force Morph tags (!!).")

    # 2. Inject Morph transitions into XML
    morph_transition_xml = (
        '<mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006">'
        '<mc:Choice xmlns:p159="http://schemas.microsoft.com/office/powerpoint/2015/09/main" Requires="p159">'
        '<p:transition spd="slow" xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" p14:dur="1250">'
        '<p159:morph option="byObject"/>'
        '</p:transition>'
        '</mc:Choice>'
        '<mc:Fallback>'
        '<p:transition spd="med">'
        '<p:fade/>'
        '</p:transition>'
        '</mc:Fallback>'
        '</mc:AlternateContent>'
    )

    temp_injected = r"C:\Users\NCN0C\Pictures\PPTMORPH\temp_injected.pptx"
    
    with zipfile.ZipFile(temp_saved, 'r') as zin, zipfile.ZipFile(temp_injected, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            # Inject to all slides (Slide 1 can have fade or morph, slides 2-12 have smooth morph)
            if item.filename.startswith('ppt/slides/slide') and item.filename.endswith('.xml'):
                xml_str = data.decode('utf-8')
                
                # Check slide index
                slide_num_match = re.search(r'slide(\d+)\.xml', item.filename)
                slide_num = int(slide_num_match.group(1)) if slide_num_match else 1
                
                # Remove any existing transition tags if present
                xml_str = re.sub(r'<p:transition.*?</p:transition>', '', xml_str)
                xml_str = re.sub(r'<mc:AlternateContent.*?</mc:AlternateContent>', '', xml_str)
                
                # We inject morph on all slides (or slide 2..12)
                # Inserting before </p:sld>
                if slide_num >= 2:
                    xml_str = xml_str.replace('</p:sld>', f'{morph_transition_xml}</p:sld>')
                else:
                    # Slide 1 fade / smooth start
                    s1_trans = '<p:transition spd="med"><p:fade/></p:transition>'
                    xml_str = xml_str.replace('</p:sld>', f'{s1_trans}</p:sld>')
                    
                data = xml_str.encode('utf-8')
            zout.writestr(item, data)

    # 3. Distribute outputs
    shutil.copy2(temp_injected, INPUT_PPTX)
    shutil.copy2(temp_injected, OUTPUT_DOWNLOADS_MORPH)
    shutil.copy2(temp_injected, OUTPUT_WORKSPACE_MORPH)
    
    # Cleanup temp files
    if os.path.exists(temp_saved):
        os.remove(temp_saved)
    if os.path.exists(temp_injected):
        os.remove(temp_injected)

    print("Morph transitions successfully injected into all slides!")
    print(f"Updated: {INPUT_PPTX}")
    print(f"Output 1: {OUTPUT_DOWNLOADS_MORPH}")
    print(f"Output 2: {OUTPUT_WORKSPACE_MORPH}")

if __name__ == "__main__":
    setup_morph_presentation()
