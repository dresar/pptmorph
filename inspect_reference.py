from pptx import Presentation

ref_path = r"C:\Users\NCN0C\Downloads\Downloads\ppt  (7).pptx.pptx"
prs = Presentation(ref_path)

print(f"Total slides: {len(prs.slides)}")
print(f"Slide Dimensions: {prs.slide_width.inches} x {prs.slide_height.inches} in ({prs.slide_width} x {prs.slide_height} EMUs)")

w = prs.slide_width
h = prs.slide_height

for idx, slide in enumerate(prs.slides):
    print(f"\n{'='*25} SLIDE {idx+1} (Total Shapes: {len(slide.shapes)}) {'='*25}")
    for s in slide.shapes:
        txt = s.text_frame.text.replace('\n', ' ') if s.has_text_frame else ''
        
        # Check in-canvas vs off-canvas
        is_in_canvas = (s.left >= 0 and s.top >= 0 and (s.left + s.width) <= w and (s.top + s.height) <= h)
        pos_tag = "IN-CANVAS" if is_in_canvas else "OFF-CANVAS"
        
        rel_x = s.left / w
        rel_y = s.top / h
        
        print(f"[{pos_tag:10s}] ID:{s.shape_id:2d} | Name:{s.name:25s} | Pos:({s.left:8d}, {s.top:8d}, {s.width:8d}, {s.height:8d}) (rel: x={rel_x:.2f}, y={rel_y:.2f}) | Text: '{txt[:40]}'")
