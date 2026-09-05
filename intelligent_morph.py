import os
import copy
import zipfile
import re
from pptx import Presentation
import xml.etree.ElementTree as ET

NS = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
}

TARGET_PPTX = r"C:\Users\NCN0C\Downloads\Desain tanpa judul.pptx"

def get_shape_transform(elem):
    off = elem.find('.//a:off', NS)
    ext = elem.find('.//a:ext', NS)
    x = int(off.get('x', 0)) if off is not None else 0
    y = int(off.get('y', 0)) if off is not None else 0
    cx = int(ext.get('cx', 0)) if ext is not None else 0
    cy = int(ext.get('cy', 0)) if ext is not None else 0
    return x, y, cx, cy

def set_shape_position(elem, new_x, new_y):
    off = elem.find('.//a:off', NS)
    if off is not None:
        off.set('x', str(int(new_x)))
        off.set('y', str(int(new_y)))

def set_shape_id_name(elem, shape_id, shape_name):
    cnvpr = elem.find('.//p:cNvPr', NS)
    if cnvpr is not None:
        cnvpr.set('id', str(shape_id))
        cnvpr.set('name', shape_name)

def generate_intelligent_morph(pptx_path):
    prs = Presentation(pptx_path)
    w = prs.slide_width
    h = prs.slide_height
    num_slides = len(prs.slides)

    # Slide Title Dictionary (H1 & H2)
    slide_titles = {
        0: ("PORTOFOLIO", "Kevin Pratama"),
        1: ("TENTANG", "SAYA"),
        2: ("PROFIL", "DIRI"),
        3: ("KEAHLIAN", "SAYA"),
        4: ("RIWAYAT", "PENDIDIKAN"),
        5: ("SERTIFIKASI", "SAYA"),
        6: ("BIDANG", "KONTEN"),
        7: ("PORTOFOLIO", "PILIHAN"),
        8: ("PENCAPAIAN", "KREATIF"),
        9: ("PENGALAMAN", "PROFESIONAL"),
        10: ("LAYANAN", "KREATIF"),
        11: ("TERIMA", "KASIH"),
    }

    # Step 1: Clean any existing cloned off-canvas shapes from previous runs
    # (Keep only shapes whose name doesn't start with !!S_cloned or identify original shapes)
    slide_original_shapes = []

    for s_idx, slide in enumerate(prs.slides):
        h1_text, h2_text = slide_titles.get(s_idx, ("", ""))
        content_shapes = []
        
        # Collect shapes
        shapes_to_keep = []
        for s in list(slide.shapes):
            txt = s.text_frame.text.replace('\n', ' ').strip() if s.has_text_frame else ''
            
            # Identify persistent anchors
            is_bg = (s.left == 0 and s.top == 0 and s.width >= w and s.height >= h)
            is_header_line = (s.top in range(900000, 1150000) and s.height < 50000 and s.width > 8000000)
            is_header_title = ("Portofolio - Kevin Pratama" in txt)
            is_page_num = (txt in [f"{i:02d}" for i in range(1, 20)] or txt in [f"{i}" for i in range(1, 20)])
            is_page_num_grp = (s.left > 14000000 and s.top < 1500000 and s.width < 2500000 and s.height < 1000000)
            is_left_decor = (s.name in ["Group 7", "!!LeftDecor", "!!Decor_Group_7"] or (s.left < 0 and s.width > 3000000 and s.height > 10000000))
            
            # Check if this shape is a title H1 or H2
            is_h1 = bool(h1_text and (txt == h1_text or txt.startswith(h1_text)))
            is_h2 = bool(h2_text and (txt == h2_text or txt.startswith(h2_text)))

            if is_bg:
                s.name = "!!Persistent_Background"
            elif is_header_line:
                s.name = "!!Persistent_HeaderLine"
            elif is_header_title:
                s.name = "!!Persistent_HeaderTitle"
            elif is_page_num or is_page_num_grp:
                s.name = "!!Persistent_PageNumber"
            elif is_left_decor:
                s.name = "!!Persistent_LeftDecor"
            elif is_h1:
                s.name = "!!Title_Word1"
            elif is_h2:
                s.name = "!!Title_Word2"
            else:
                # Content shape (Card, Body text, Icon, Stats, Timeline)
                morph_name = f"!!S{s_idx+1}_Shp{s.shape_id}"
                s.name = morph_name
                
                elem_copy = copy.deepcopy(s._element)
                x, y, cx, cy = get_shape_transform(elem_copy)
                
                # Check if it was originally in-canvas
                in_canvas = (s.left >= -500000 and s.top >= -500000 and s.left + s.width <= w + 1000000 and s.top + s.height <= h + 1000000)
                if in_canvas:
                    content_shapes.append({
                        'elem': elem_copy,
                        'orig_x': x,
                        'orig_y': y,
                        'cx': cx,
                        'cy': cy,
                        'name': morph_name,
                        'text': txt
                    })

        slide_original_shapes.append(content_shapes)

    temp_base = pptx_path + ".base.tmp.pptx"
    prs.save(temp_base)

    # Step 2: Inject Multi-Vector Staggered Off-Canvas Staging
    prs = Presentation(temp_base)

    def get_next_id(slide):
        max_id = 200
        for sh in slide.shapes:
            if sh.shape_id > max_id:
                max_id = sh.shape_id
        return max_id + 20

    for i in range(num_slides):
        curr_slide = prs.slides[i]
        sp_tree = curr_slide._element.spTree
        next_id = get_next_id(curr_slide)

        # INBOUND (Upcoming from Slide i+1 staged in Slide i)
        if i + 1 < num_slides:
            next_items = slide_original_shapes[i + 1]
            for idx_item, item in enumerate(next_items):
                elem = copy.deepcopy(item['elem'])
                orig_x = item['orig_x']
                orig_y = item['orig_y']
                stagger_offset = (idx_item % 4) * 350000

                # Intelligent multi-vector entry:
                # Cards at bottom half slide in from bottom
                # Elements at right half slide in from right
                if orig_y > h * 0.45:
                    entry_x = orig_x
                    entry_y = int(h + 2000000 + stagger_offset)
                else:
                    entry_x = int(w + 2500000 + stagger_offset)
                    entry_y = orig_y

                set_shape_position(elem, entry_x, entry_y)
                set_shape_id_name(elem, next_id, item['name'])
                next_id += 1
                sp_tree.append(elem)

        # OUTBOUND (Departing from Slide i-1 staged in Slide i)
        if i > 0:
            prev_items = slide_original_shapes[i - 1]
            for idx_item, item in enumerate(prev_items):
                elem = copy.deepcopy(item['elem'])
                orig_x = item['orig_x']
                orig_y = item['orig_y']
                cx = item['cx']
                cy = item['cy']
                stagger_offset = (idx_item % 4) * 350000

                # Intelligent multi-vector exit:
                # Top elements exit to top, lower elements exit to left
                if orig_y < h * 0.4:
                    exit_x = orig_x
                    exit_y = int(-cy - 2000000 - stagger_offset)
                else:
                    exit_x = int(-cx - 2500000 - stagger_offset)
                    exit_y = orig_y

                set_shape_position(elem, exit_x, exit_y)
                set_shape_id_name(elem, next_id, item['name'])
                next_id += 1
                sp_tree.append(elem)

    temp_staged = pptx_path + ".staged.tmp.pptx"
    prs.save(temp_staged)

    # Step 3: Inject Modern XML Morph Transition
    morph_xml = (
        '<mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006">'
        '<mc:Choice xmlns:p159="http://schemas.microsoft.com/office/powerpoint/2015/09/main" Requires="p159">'
        '<p:transition spd="slow" xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" p14:dur="1250">'
        '<p159:morph option="byObject"/>'
        '</p:transition>'
        '</mc:Choice>'
        '<mc:Fallback><p:transition spd="med"><p:fade/></p:transition></mc:Fallback>'
        '</mc:AlternateContent>'
    )

    temp_output = pptx_path + ".out.tmp.pptx"
    with zipfile.ZipFile(temp_staged, 'r') as zin, zipfile.ZipFile(temp_output, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith('ppt/slides/slide') and item.filename.endswith('.xml'):
                xml_str = data.decode('utf-8')
                slide_num_match = re.search(r'slide(\d+)\.xml', item.filename)
                slide_num = int(slide_num_match.group(1)) if slide_num_match else 1
                
                xml_str = re.sub(r'<p:transition.*?</p:transition>', '', xml_str)
                xml_str = re.sub(r'<mc:AlternateContent.*?</mc:AlternateContent>', '', xml_str)
                
                if slide_num >= 2:
                    xml_str = xml_str.replace('</p:sld>', f'{morph_xml}</p:sld>')
                else:
                    xml_str = xml_str.replace('</p:sld>', '<p:transition spd="med"><p:fade/></p:transition></p:sld>')
                data = xml_str.encode('utf-8')
            zout.writestr(item, data)

    # Directly overwrite target file
    if os.path.exists(TARGET_PPTX):
        os.remove(TARGET_PPTX)
    os.rename(temp_output, TARGET_PPTX)

    # Clean temporary files
    for tmp in [temp_base, temp_staged]:
        if os.path.exists(tmp):
            os.remove(tmp)

    print(f"Intelligent Off-Canvas Morph successfully built directly into: {TARGET_PPTX}")

if __name__ == "__main__":
    generate_intelligent_morph(TARGET_PPTX)
