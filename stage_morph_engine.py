import os
import shutil
import copy
import zipfile
import re
from pptx import Presentation
import xml.etree.ElementTree as ET

NS = {
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
}

INPUT_PPTX = r"C:\Users\NCN0C\Downloads\Desain tanpa judul.pptx"
BACKUP_PPTX = r"C:\Users\NCN0C\Downloads\Desain tanpa judul_backup.pptx"
OUTPUT_DOWNLOADS = r"C:\Users\NCN0C\Downloads\Desain tanpa judul_morph.pptx"
OUTPUT_WORKSPACE = r"C:\Users\NCN0C\Pictures\PPTMORPH\Desain tanpa judul_morph.pptx"

def get_shape_transform(elem):
    off_elem = elem.find('.//a:off', NS)
    ext_elem = elem.find('.//a:ext', NS)
    x = int(off_elem.get('x', 0)) if off_elem is not None else 0
    y = int(off_elem.get('y', 0)) if off_elem is not None else 0
    cx = int(ext_elem.get('cx', 0)) if ext_elem is not None else 0
    cy = int(ext_elem.get('cy', 0)) if ext_elem is not None else 0
    return x, y, cx, cy

def set_shape_position(elem, new_x, new_y):
    off_elem = elem.find('.//a:off', NS)
    if off_elem is not None:
        off_elem.set('x', str(int(new_x)))
        off_elem.set('y', str(int(new_y)))

def set_shape_id_name(elem, shape_id, shape_name):
    # Check cNvPr in nvSpPr, nvGrpSpPr, nvPicPr, nvGraphicFramePr
    cnvpr = elem.find('.//p:cNvPr', NS)
    if cnvpr is not None:
        cnvpr.set('id', str(shape_id))
        cnvpr.set('name', shape_name)

def build_off_canvas_morph():
    # 1. Ensure backup
    if not os.path.exists(BACKUP_PPTX):
        shutil.copy2(INPUT_PPTX, BACKUP_PPTX)

    prs = Presentation(BACKUP_PPTX)
    w = prs.slide_width
    h = prs.slide_height
    num_slides = len(prs.slides)

    # 2. Tag and classify shapes in each slide
    slide_content_shapes = [] # list of lists of (shape_elem, original_pos, name)
    
    for s_idx, slide in enumerate(prs.slides):
        content_for_this_slide = []
        
        for shape in slide.shapes:
            txt = shape.text_frame.text.replace('\n', ' ').strip() if shape.has_text_frame else ''
            
            # Persistent checks
            is_bg = (shape.left == 0 and shape.top == 0 and shape.width >= w and shape.height >= h)
            is_header_line = (shape.top in range(900000, 1150000) and shape.height < 50000 and shape.width > 8000000)
            is_header_title = ("Portofolio - Kevin Pratama" in txt)
            is_page_num = (txt in [f"{i:02d}" for i in range(1, 20)] or txt in [f"{i}" for i in range(1, 20)])
            is_page_num_grp = (shape.left > 14000000 and shape.top < 1500000 and shape.width < 2500000 and shape.height < 1000000)
            
            if is_bg:
                shape.name = "!!Persistent_Background"
            elif is_header_line:
                shape.name = "!!Persistent_HeaderLine"
            elif is_header_title:
                shape.name = "!!Persistent_HeaderTitle"
            elif is_page_num or is_page_num_grp:
                shape.name = "!!Persistent_PageNumber"
            else:
                # This is a transient content shape!
                # Give it a unique morph name tied to this slide & shape ID
                morph_name = f"!!S{s_idx+1}_Shp{shape.shape_id}"
                shape.name = morph_name
                
                elem_copy = copy.deepcopy(shape._element)
                x, y, cx, cy = get_shape_transform(elem_copy)
                content_for_this_slide.append({
                    'elem': elem_copy,
                    'orig_x': x,
                    'orig_y': y,
                    'cx': cx,
                    'cy': cy,
                    'name': morph_name,
                    'text': txt
                })
                
        slide_content_shapes.append(content_for_this_slide)

    # Save tagged base
    temp_tagged = r"C:\Users\NCN0C\Pictures\PPTMORPH\temp_tagged.pptx"
    prs.save(temp_tagged)

    # 3. Reload tagged presentation to perform Off-Canvas Cloning
    prs = Presentation(temp_tagged)
    
    # ID counter generator per slide
    def get_max_id(slide):
        max_id = 100
        for sh in slide.shapes:
            if sh.shape_id > max_id:
                max_id = sh.shape_id
        return max_id + 50

    for i in range(num_slides):
        curr_slide = prs.slides[i]
        curr_sp_tree = curr_slide._element.spTree
        next_id = get_max_id(curr_slide)

        # A. INBOUND STAGING: Add Slide i+1's shapes to Slide i (Off-Canvas Right / Bottom)
        if i + 1 < num_slides:
            next_slide_shapes = slide_content_shapes[i + 1]
            for item in next_slide_shapes:
                inbound_elem = copy.deepcopy(item['elem'])
                
                # Determine off-canvas entry vector
                # Enter from right or bottom depending on position
                orig_x = item['orig_x']
                orig_y = item['orig_y']
                
                # If element is near top, enter from top/right, if near bottom enter from bottom
                if orig_y < h * 0.4:
                    # Enter from Right-Top
                    entry_x = int(w + 3000000 + (orig_x % 1000000))
                    entry_y = orig_y
                else:
                    # Enter from Right-Bottom
                    entry_x = int(orig_x + w * 0.8)
                    entry_y = int(h + 2500000 + (orig_y % 1000000))
                
                set_shape_position(inbound_elem, entry_x, entry_y)
                set_shape_id_name(inbound_elem, next_id, item['name'])
                next_id += 1
                
                # Append to current slide spTree
                curr_sp_tree.append(inbound_elem)

        # B. OUTBOUND STAGING: Add Slide i-1's shapes to Slide i (Off-Canvas Left / Top)
        if i > 0:
            prev_slide_shapes = slide_content_shapes[i - 1]
            for item in prev_slide_shapes:
                outbound_elem = copy.deepcopy(item['elem'])
                
                orig_x = item['orig_x']
                orig_y = item['orig_y']
                
                # Exit to Left / Top
                exit_x = int(-item['cx'] - 4000000 - (orig_x % 1000000))
                if orig_y < h * 0.5:
                    exit_y = int(-item['cy'] - 2000000)
                else:
                    exit_y = orig_y
                
                set_shape_position(outbound_elem, exit_x, exit_y)
                set_shape_id_name(outbound_elem, next_id, item['name'])
                next_id += 1
                
                # Append to current slide spTree
                curr_sp_tree.append(outbound_elem)

    temp_staged = r"C:\Users\NCN0C\Pictures\PPTMORPH\temp_staged.pptx"
    prs.save(temp_staged)
    print("Off-Canvas Staging elements injected into all slides successfully!")

    # 4. Inject Morph Transition XML
    morph_transition_xml = (
        '<mc:AlternateContent xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006">'
        '<mc:Choice xmlns:p159="http://schemas.microsoft.com/office/powerpoint/2015/09/main" Requires="p159">'
        '<p:transition spd="slow" xmlns:p14="http://schemas.microsoft.com/office/powerpoint/2010/main" p14:dur="1250">'
        '<p159:morph option="byObject"/>'
        '</p:transition>'
        '</mc:Choice>'
        '<mc:Fallback>'
        '<p:transition spd="med">'
        '<p:fade/>'
        '</p:transition>'
        '</mc:Fallback>'
        '</mc:AlternateContent>'
    )

    temp_final = r"C:\Users\NCN0C\Pictures\PPTMORPH\temp_final.pptx"
    with zipfile.ZipFile(temp_staged, 'r') as zin, zipfile.ZipFile(temp_final, 'w', zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith('ppt/slides/slide') and item.filename.endswith('.xml'):
                xml_str = data.decode('utf-8')
                
                slide_num_match = re.search(r'slide(\d+)\.xml', item.filename)
                slide_num = int(slide_num_match.group(1)) if slide_num_match else 1
                
                xml_str = re.sub(r'<p:transition.*?</p:transition>', '', xml_str)
                xml_str = re.sub(r'<mc:AlternateContent.*?</mc:AlternateContent>', '', xml_str)
                
                if slide_num >= 2:
                    xml_str = xml_str.replace('</p:sld>', f'{morph_transition_xml}</p:sld>')
                else:
                    s1_trans = '<p:transition spd="med"><p:fade/></p:transition>'
                    xml_str = xml_str.replace('</p:sld>', f'{s1_trans}</p:sld>')
                    
                data = xml_str.encode('utf-8')
            zout.writestr(item, data)

    # Distribute outputs safely
    saved_files = []
    
    # 1. Workspace copy (always accessible)
    shutil.copy2(temp_final, OUTPUT_WORKSPACE)
    saved_files.append(OUTPUT_WORKSPACE)
    
    # 2. Main target in Downloads
    try:
        shutil.copy2(temp_final, INPUT_PPTX)
        saved_files.append(INPUT_PPTX)
    except PermissionError:
        fallback_input = r"C:\Users\NCN0C\Downloads\Desain tanpa judul_staged_morph.pptx"
        shutil.copy2(temp_final, fallback_input)
        saved_files.append(fallback_input)
        print(f"Note: {INPUT_PPTX} is open in PowerPoint. Saved as {fallback_input}")

    # 3. Morph copy in Downloads
    try:
        shutil.copy2(temp_final, OUTPUT_DOWNLOADS)
        saved_files.append(OUTPUT_DOWNLOADS)
    except PermissionError:
        fallback_morph = r"C:\Users\NCN0C\Downloads\Desain tanpa judul_morph_v2.pptx"
        shutil.copy2(temp_final, fallback_morph)
        saved_files.append(fallback_morph)
        print(f"Note: {OUTPUT_DOWNLOADS} is open in PowerPoint. Saved as {fallback_morph}")

    # Clean temporary files
    for tmp in [temp_tagged, temp_staged, temp_final]:
        if os.path.exists(tmp):
            os.remove(tmp)

    print("Complete Off-Canvas Morph Master Presentation generated successfully!")
    for f in saved_files:
        print(f"Saved: {f}")

if __name__ == "__main__":
    build_off_canvas_morph()
